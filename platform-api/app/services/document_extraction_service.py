"""Extract text from unstructured documents (PDF, DOCX, PPTX, TXT, Markdown)."""

from __future__ import annotations

import logging
from pathlib import Path
from typing import Any

logger = logging.getLogger(__name__)


def extract_text(file_path: str, file_extension: str) -> dict[str, Any]:
    """Return structured extraction: document_text, sections[], metadata."""
    ext = file_extension.lower()
    if ext == ".pdf":
        return _extract_pdf(file_path)
    elif ext == ".docx":
        return _extract_docx(file_path)
    elif ext == ".pptx":
        return _extract_pptx(file_path)
    elif ext in (".txt", ".md"):
        return _extract_text_file(file_path)
    else:
        raise ValueError(f"Unsupported file type: {ext}")


def _extract_pdf(file_path: str) -> dict[str, Any]:
    try:
        from pypdf import PdfReader
    except ImportError:
        logger.warning("pypdf not installed, falling back to raw read")
        return _extract_text_file(file_path)

    reader = PdfReader(file_path)
    sections = []
    all_text_parts = []
    for i, page in enumerate(reader.pages, 1):
        text = page.extract_text() or ""
        sections.append({
            "section_type": "page",
            "section_number": i,
            "title": None,
            "text": text,
        })
        all_text_parts.append(text)

    full_text = "\n\n".join(all_text_parts)
    return {
        "document_text": full_text,
        "sections": sections,
        "metadata": {
            "page_count": len(reader.pages),
            "slide_count": None,
            "word_count": len(full_text.split()),
        },
    }


def _extract_docx(file_path: str) -> dict[str, Any]:
    try:
        from docx import Document
    except ImportError:
        logger.warning("python-docx not installed, falling back to raw read")
        return _extract_text_file(file_path)

    doc = Document(file_path)
    sections = []
    all_text_parts = []
    section_num = 0

    for para in doc.paragraphs:
        if not para.text.strip():
            continue
        section_num += 1
        is_heading = para.style and para.style.name and para.style.name.startswith("Heading")
        sections.append({
            "section_type": "heading" if is_heading else "paragraph",
            "section_number": section_num,
            "title": para.text if is_heading else None,
            "text": para.text,
        })
        all_text_parts.append(para.text)

    # Extract tables
    for table in doc.tables:
        section_num += 1
        rows_text = []
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            rows_text.append(" | ".join(cells))
        table_text = "Table:\n" + "\n".join(rows_text)
        sections.append({
            "section_type": "table",
            "section_number": section_num,
            "title": None,
            "text": table_text,
        })
        all_text_parts.append(table_text)

    full_text = "\n\n".join(all_text_parts)
    return {
        "document_text": full_text,
        "sections": sections,
        "metadata": {
            "page_count": None,
            "slide_count": None,
            "word_count": len(full_text.split()),
        },
    }


def _extract_pptx(file_path: str) -> dict[str, Any]:
    try:
        from pptx import Presentation
    except ImportError:
        logger.warning("python-pptx not installed, falling back to raw read")
        return _extract_text_file(file_path)

    prs = Presentation(file_path)
    sections = []
    all_text_parts = []

    for i, slide in enumerate(prs.slides, 1):
        slide_texts = []
        title = None
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text
                if shape == slide.shapes.title:
                    title = text
                slide_texts.append(text)
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    slide_texts.append(" | ".join(cells))

        slide_text = "\n".join(slide_texts)
        sections.append({
            "section_type": "slide",
            "section_number": i,
            "title": title,
            "text": slide_text,
        })
        all_text_parts.append(slide_text)

    full_text = "\n\n".join(all_text_parts)
    return {
        "document_text": full_text,
        "sections": sections,
        "metadata": {
            "page_count": None,
            "slide_count": len(prs.slides),
            "word_count": len(full_text.split()),
        },
    }


def _extract_text_file(file_path: str) -> dict[str, Any]:
    text = Path(file_path).read_text(encoding="utf-8", errors="replace")
    return {
        "document_text": text,
        "sections": [{
            "section_type": "text",
            "section_number": 1,
            "title": None,
            "text": text,
        }],
        "metadata": {
            "page_count": None,
            "slide_count": None,
            "word_count": len(text.split()),
        },
    }
