"""Server-side structured previews for project_assets documents.

Every format is parsed locally with an already-vendored library (or the
local `antiword` binary for legacy .doc) and reduced to a small, size-
bounded JSON shape the frontend viewer renders directly. Uploaded documents
are never sent to an external preview service (Google Docs Viewer, Office
Viewer, etc.) -- doing so would leak tenant-private content and break the
project/tenant isolation the rest of this feature enforces.
"""

from __future__ import annotations

import io
import logging
import shutil
import subprocess
import tempfile
import xml.dom.minidom
from collections.abc import Callable
from typing import Any

logger = logging.getLogger(__name__)

# Keeps the viewer responsive against a huge document instead of shipping
# (or fully parsing) the whole thing into one HTTP response.
MAX_TEXT_CHARS = 250_000
MAX_SHEET_ROWS = 500
MAX_SHEET_COLS = 50
MAX_SHEETS = 20
MAX_SECTIONS = 500
MAX_PREVIEW_FILE_BYTES = 25 * 1024 * 1024  # 25 MB

# Formats the browser renders natively -- the viewer fetches these as an
# authenticated blob and hands them straight to <embed>/<img> rather than
# routing them through this module at all.
NATIVE_EXTENSIONS = frozenset({
    ".pdf", ".jpg", ".jpeg", ".png", ".gif", ".webp", ".bmp", ".svg",
})

TEXT_EXTENSIONS = frozenset({
    ".txt", ".md", ".csv", ".tsv", ".json", ".xml", ".yaml", ".yml", ".log",
})


def _truncate_text(text: str) -> tuple[str, bool]:
    if len(text) > MAX_TEXT_CHARS:
        return text[:MAX_TEXT_CHARS], True
    return text, False


def _json_safe(value: Any) -> Any:
    if value is None or isinstance(value, str | int | float | bool):
        return value
    return str(value)


def _preview_text(data: bytes, ext: str) -> dict[str, Any]:
    text = data.decode("utf-8", errors="replace")
    if ext == ".xml":
        try:
            text = xml.dom.minidom.parseString(text).toprettyxml(indent="  ")
        except Exception:
            pass  # Malformed XML still previews fine as raw text.
    body, truncated = _truncate_text(text)
    return {"kind": "text", "text": body, "truncated": truncated}


def _preview_docx(data: bytes) -> dict[str, Any]:
    from docx import Document

    document = Document(io.BytesIO(data))
    paragraphs = [p.text for p in document.paragraphs if p.text.strip()]
    return {
        "kind": "docx",
        "paragraphs": paragraphs[:MAX_SECTIONS],
        "truncated": len(paragraphs) > MAX_SECTIONS,
    }


def _preview_pptx(data: bytes) -> dict[str, Any]:
    from pptx import Presentation

    presentation = Presentation(io.BytesIO(data))
    slides = []
    for index, slide in enumerate(presentation.slides):
        if index >= MAX_SECTIONS:
            break
        texts = [
            shape.text_frame.text
            for shape in slide.shapes
            if shape.has_text_frame and shape.text_frame.text.strip()
        ]
        slides.append({"index": index + 1, "texts": texts})
    return {
        "kind": "pptx",
        "slides": slides,
        "truncated": len(presentation.slides) > MAX_SECTIONS,
    }


def _preview_xlsx(data: bytes) -> dict[str, Any]:
    from openpyxl import load_workbook

    workbook = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    try:
        sheet_names = workbook.sheetnames
        sheets = []
        for name in sheet_names[:MAX_SHEETS]:
            worksheet = workbook[name]
            rows: list[list[Any]] = []
            total_rows = 0
            total_cols = 0
            for row_index, row in enumerate(worksheet.iter_rows(values_only=True)):
                total_rows = row_index + 1
                total_cols = max(total_cols, len(row))
                if row_index < MAX_SHEET_ROWS:
                    rows.append([_json_safe(v) for v in row[:MAX_SHEET_COLS]])
            sheets.append({
                "name": name,
                "rows": rows,
                "totalRows": total_rows,
                "totalCols": total_cols,
                "truncatedRows": total_rows > MAX_SHEET_ROWS,
                "truncatedCols": total_cols > MAX_SHEET_COLS,
            })
        return {
            "kind": "spreadsheet",
            "sheets": sheets,
            "truncatedSheets": len(sheet_names) > MAX_SHEETS,
        }
    finally:
        workbook.close()


def _preview_xls(data: bytes) -> dict[str, Any]:
    import xlrd

    workbook = xlrd.open_workbook(file_contents=data)
    sheets = []
    for sheet in workbook.sheets()[:MAX_SHEETS]:
        rows = [
            [_json_safe(v) for v in sheet.row_values(row_index)[:MAX_SHEET_COLS]]
            for row_index in range(min(sheet.nrows, MAX_SHEET_ROWS))
        ]
        sheets.append({
            "name": sheet.name,
            "rows": rows,
            "totalRows": sheet.nrows,
            "totalCols": sheet.ncols,
            "truncatedRows": sheet.nrows > MAX_SHEET_ROWS,
            "truncatedCols": sheet.ncols > MAX_SHEET_COLS,
        })
    return {
        "kind": "spreadsheet",
        "sheets": sheets,
        "truncatedSheets": workbook.nsheets > MAX_SHEETS,
    }


def _preview_doc(data: bytes) -> dict[str, Any]:
    """Legacy binary .doc via the local `antiword` CLI. Falls back to
    "unsupported" (download-only) rather than raising when the binary isn't
    installed or the file can't be converted -- a bad/corrupt upload must
    never break the viewer, only degrade to a download link."""
    if shutil.which("antiword") is None:
        return {"kind": "unsupported", "reason": "Preview is unavailable for this file."}
    with tempfile.NamedTemporaryFile(suffix=".doc") as tmp:
        tmp.write(data)
        tmp.flush()
        try:
            result = subprocess.run(
                ["antiword", tmp.name],
                capture_output=True, timeout=20, check=False,
            )
        except (OSError, subprocess.TimeoutExpired):
            return {"kind": "unsupported", "reason": "This file could not be converted for preview."}
    if result.returncode != 0 or not result.stdout:
        return {"kind": "unsupported", "reason": "This file could not be converted for preview."}
    body, truncated = _truncate_text(result.stdout.decode("utf-8", errors="replace"))
    return {"kind": "text", "text": body, "truncated": truncated}


_BUILDERS: dict[str, Callable[[bytes], dict[str, Any]]] = {
    ".docx": _preview_docx,
    ".pptx": _preview_pptx,
    ".xlsx": _preview_xlsx,
    ".xlsm": _preview_xlsx,
    ".xls": _preview_xls,
    ".doc": _preview_doc,
}


def build_preview(
    *,
    file_extension: str | None,
    file_size_bytes: int | None,
    read_bytes: Callable[[], bytes],
) -> dict[str, Any]:
    """Build a bounded JSON preview for one document.

    ``read_bytes`` is only called when a format actually needs the file's
    bytes -- a native (PDF/image) or oversized file returns without reading
    the file at all.
    """
    ext = (file_extension or "").lower()
    if ext in NATIVE_EXTENSIONS:
        return {"kind": "native"}
    if file_size_bytes is not None and file_size_bytes > MAX_PREVIEW_FILE_BYTES:
        return {"kind": "unsupported", "reason": "File is too large to preview -- use Download instead."}

    if ext in TEXT_EXTENSIONS:
        try:
            return _preview_text(read_bytes(), ext)
        except Exception:
            logger.exception("Text preview failed for extension %s", ext)
            return {"kind": "unsupported", "reason": "This file could not be previewed."}

    builder = _BUILDERS.get(ext)
    if builder is None:
        return {"kind": "unsupported", "reason": f"Preview is not available for {ext or 'this file type'}."}
    try:
        return builder(read_bytes())
    except Exception:
        logger.exception("Preview build failed for extension %s", ext)
        return {"kind": "unsupported", "reason": "This file could not be previewed."}
