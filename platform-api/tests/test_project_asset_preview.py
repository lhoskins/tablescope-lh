"""The authenticated project document viewer: bounded structured previews
(app.services.document_preview) and the /preview + /content endpoints'
tenant, project-membership and private-document authorization.

Run from ``platform-api``: ``pytest -q tests/test_project_asset_preview.py``.
"""

from __future__ import annotations

import io

import pytest

from app.auth.jwt import create_access_token
from app.services import document_preview
from app.services.supabase_auth_service import SupabaseAuthService, SupabaseUser

pytestmark = pytest.mark.anyio


class _FakeSupabase(SupabaseAuthService):
    def __init__(self) -> None:
        pass

    async def create_or_invite_user(
        self, email, *, first_name=None, last_name=None, redirect_to=None
    ) -> SupabaseUser:
        return SupabaseUser(id=f"supa-{email}", email=email, created=True, action_link=f"https://invite/{email}")


class _FakeEmail:
    async def send_transactional_email(self, *, to, template, variables, subject=None, reply_to=None) -> bool:
        return True


@pytest.fixture(autouse=True)
def _mock_supabase(monkeypatch):
    import app.routes.tenants_users as tenants_module

    monkeypatch.setattr(tenants_module, "SupabaseAuthService", _FakeSupabase)
    monkeypatch.setattr(tenants_module, "EmailService", _FakeEmail)


# ── document_preview: pure parsing, no DB/HTTP needed ──────────────────

def test_native_extensions_never_read_the_file():
    read = lambda: (_ for _ in ()).throw(AssertionError("must not read a native file's bytes"))  # noqa: E731
    for ext in (".pdf", ".jpg", ".png", ".svg"):
        assert document_preview.build_preview(file_extension=ext, file_size_bytes=999, read_bytes=read) == {"kind": "native"}


def test_oversized_file_is_rejected_before_reading():
    read = lambda: (_ for _ in ()).throw(AssertionError("must not read an oversized file"))  # noqa: E731
    result = document_preview.build_preview(
        file_extension=".docx",
        file_size_bytes=document_preview.MAX_PREVIEW_FILE_BYTES + 1,
        read_bytes=read,
    )
    assert result["kind"] == "unsupported"


def test_text_preview_is_truncated_at_the_character_limit():
    body = "x" * (document_preview.MAX_TEXT_CHARS + 500)
    result = document_preview.build_preview(
        file_extension=".txt", file_size_bytes=len(body), read_bytes=lambda: body.encode(),
    )
    assert result["kind"] == "text"
    assert result["truncated"] is True
    assert len(result["text"]) == document_preview.MAX_TEXT_CHARS


def test_docx_preview_extracts_paragraphs():
    from docx import Document

    document = Document()
    document.add_paragraph("First paragraph")
    document.add_paragraph("Second paragraph")
    buf = io.BytesIO()
    document.save(buf)
    data = buf.getvalue()

    result = document_preview.build_preview(file_extension=".docx", file_size_bytes=len(data), read_bytes=lambda: data)
    assert result == {"kind": "docx", "paragraphs": ["First paragraph", "Second paragraph"], "truncated": False}


def test_pptx_preview_extracts_slide_text():
    from pptx import Presentation

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "Quarterly Review"
    buf = io.BytesIO()
    presentation.save(buf)
    data = buf.getvalue()

    result = document_preview.build_preview(file_extension=".pptx", file_size_bytes=len(data), read_bytes=lambda: data)
    assert result["kind"] == "pptx"
    assert result["slides"][0]["texts"] == ["Quarterly Review"]


def test_xlsx_preview_extracts_sheet_rows_and_bounds_are_reported():
    from openpyxl import Workbook

    workbook = Workbook()
    worksheet = workbook.active
    worksheet.append(["name", "amount"])
    worksheet.append(["Acme", 100])
    buf = io.BytesIO()
    workbook.save(buf)
    data = buf.getvalue()

    result = document_preview.build_preview(file_extension=".xlsx", file_size_bytes=len(data), read_bytes=lambda: data)
    assert result["kind"] == "spreadsheet"
    sheet = result["sheets"][0]
    assert sheet["rows"] == [["name", "amount"], ["Acme", 100]]
    assert sheet["totalRows"] == 2
    assert sheet["truncatedRows"] is False


def test_xlsx_preview_bounds_rows_columns_and_sheets():
    from openpyxl import Workbook

    workbook = Workbook()
    worksheet = workbook.active
    for row in range(document_preview.MAX_SHEET_ROWS + 10):
        worksheet.append([row] * (document_preview.MAX_SHEET_COLS + 5))
    buf = io.BytesIO()
    workbook.save(buf)
    data = buf.getvalue()

    result = document_preview.build_preview(file_extension=".xlsx", file_size_bytes=len(data), read_bytes=lambda: data)
    sheet = result["sheets"][0]
    assert len(sheet["rows"]) == document_preview.MAX_SHEET_ROWS
    assert len(sheet["rows"][0]) == document_preview.MAX_SHEET_COLS
    assert sheet["truncatedRows"] is True
    assert sheet["truncatedCols"] is True


def test_legacy_xls_preview():
    import xlwt

    workbook = xlwt.Workbook()
    sheet = workbook.add_sheet("Data")
    sheet.write(0, 0, "name")
    sheet.write(1, 0, "Acme")
    buf = io.BytesIO()
    workbook.save(buf)
    data = buf.getvalue()

    result = document_preview.build_preview(file_extension=".xls", file_size_bytes=len(data), read_bytes=lambda: data)
    assert result["kind"] == "spreadsheet"
    assert result["sheets"][0]["name"] == "Data"
    assert result["sheets"][0]["rows"][0][0] == "name"


def test_legacy_doc_falls_back_to_unsupported_on_bad_input():
    # A corrupt/garbage .doc must degrade to "unsupported" (download-only),
    # never raise -- a bad upload must not break the viewer.
    result = document_preview.build_preview(
        file_extension=".doc", file_size_bytes=20, read_bytes=lambda: b"not a real doc file",
    )
    assert result["kind"] == "unsupported"


def test_unknown_extension_is_unsupported():
    result = document_preview.build_preview(file_extension=".exe", file_size_bytes=10, read_bytes=lambda: b"MZ")
    assert result["kind"] == "unsupported"


def test_xml_preview_is_pretty_printed_and_falls_back_on_malformed_input():
    data = b"<root><a>1</a></root>"
    result = document_preview.build_preview(file_extension=".xml", file_size_bytes=len(data), read_bytes=lambda: data)
    assert result["kind"] == "text"
    assert "<a>1</a>" in result["text"] or "<a>\n" in result["text"]

    malformed = b"<root><a>1</a>"
    result2 = document_preview.build_preview(file_extension=".xml", file_size_bytes=len(malformed), read_bytes=lambda: malformed)
    assert result2["kind"] == "text"  # falls back to raw text instead of raising


# ── /preview + /content: auth and containment ───────────────────────────

def _headers(tenant_id: int, user_id: int, role: str = "viewer") -> dict:
    token = create_access_token(sub="u", tenant_id=tenant_id, user_id=user_id, role=role)
    return {"Authorization": f"Bearer {token}"}


async def _make_tenant_user(client, service_headers, slug: str) -> tuple[dict, dict]:
    r = await client.post("/api/tenants", json={"slug": slug, "name": slug}, headers=service_headers)
    assert r.status_code == 201
    tenant = r.json()
    r = await client.post(
        f"/api/tenants/{tenant['id']}/users",
        json={"email": f"{slug}@test.com", "display_name": "U", "role": "editor", "external_id": f"ext-{slug}"},
        headers=service_headers,
    )
    assert r.status_code == 201
    user = r.json()
    return tenant, user


async def _seed_asset(
    db_session, monkeypatch, tmp_path, *,
    tenant_id: int, project_id: int, owner_user_id: int | None,
    visibility: str = "shared_project", content: bytes = b"hello world",
    extension: str = ".txt",
):
    import app.routes.project_assets as project_assets_module
    from app.models.project_asset import ProjectAsset

    monkeypatch.setattr(project_assets_module, "LOCAL_STORAGE_BASE", str(tmp_path))
    file_path = tmp_path / f"asset{extension}"
    file_path.write_bytes(content)

    asset = ProjectAsset(
        tenant_id=tenant_id, project_id=project_id, owner_user_id=owner_user_id,
        asset_type="txt", source_type="uploaded_file", title="Doc",
        filename=f"asset{extension}", original_filename=f"asset{extension}",
        content_type="text/plain", file_extension=extension,
        storage_provider="local", storage_location=str(file_path),
        file_size_bytes=len(content), visibility=visibility, status="uploaded",
    )
    db_session.add(asset)
    await db_session.commit()
    await db_session.refresh(asset)
    return asset


async def test_project_member_can_preview_and_download_a_shared_asset(client, service_headers, db_session, monkeypatch, tmp_path):
    tenant, owner = await _make_tenant_user(client, service_headers, "doc-view-member")
    r = await client.post(
        "/api/projects", json={"name": "P", "description": "", "is_shared": True},
        headers=_headers(tenant["id"], owner["id"], "editor"),
    )
    assert r.status_code == 201
    project = r.json()

    asset = await _seed_asset(
        db_session, monkeypatch, tmp_path,
        tenant_id=tenant["id"], project_id=project["id"], owner_user_id=owner["id"],
        content=b"Hello, viewer!", extension=".txt",
    )

    r = await client.get(
        f"/api/projects/{project['id']}/assets/{asset.id}/preview",
        headers=_headers(tenant["id"], owner["id"]),
    )
    assert r.status_code == 200
    body = r.json()
    assert body["kind"] == "text"
    assert body["text"] == "Hello, viewer!"
    assert r.headers.get("cache-control") == "private, no-store"
    assert r.headers.get("x-content-type-options") == "nosniff"

    r = await client.get(
        f"/api/projects/{project['id']}/assets/{asset.id}/content",
        headers=_headers(tenant["id"], owner["id"]),
    )
    assert r.status_code == 200
    assert r.content == b"Hello, viewer!"
    assert r.headers.get("cache-control") == "private, no-store"
    assert r.headers.get("x-content-type-options") == "nosniff"


async def test_non_member_cannot_preview_or_download_from_a_shared_project(client, service_headers, db_session, monkeypatch, tmp_path):
    tenant, owner = await _make_tenant_user(client, service_headers, "doc-view-nonmember")
    r = await client.post(
        "/api/projects", json={"name": "P", "description": "", "is_shared": True},
        headers=_headers(tenant["id"], owner["id"], "editor"),
    )
    project = r.json()
    asset = await _seed_asset(
        db_session, monkeypatch, tmp_path,
        tenant_id=tenant["id"], project_id=project["id"], owner_user_id=owner["id"],
    )

    r = await client.post(
        f"/api/tenants/{tenant['id']}/users",
        json={"email": "outsider@test.com", "display_name": "Outsider", "role": "viewer", "external_id": "ext-outsider"},
        headers=service_headers,
    )
    outsider = r.json()

    r = await client.get(
        f"/api/projects/{project['id']}/assets/{asset.id}/preview",
        headers=_headers(tenant["id"], outsider["id"]),
    )
    assert r.status_code == 403
    r = await client.get(
        f"/api/projects/{project['id']}/assets/{asset.id}/content",
        headers=_headers(tenant["id"], outsider["id"]),
    )
    assert r.status_code == 403


async def test_another_tenant_cannot_access_the_asset_even_knowing_both_ids(client, service_headers, db_session, monkeypatch, tmp_path):
    tenant, owner = await _make_tenant_user(client, service_headers, "doc-view-tenant-a")
    r = await client.post(
        "/api/projects", json={"name": "P", "description": "", "is_shared": True},
        headers=_headers(tenant["id"], owner["id"], "editor"),
    )
    project = r.json()
    asset = await _seed_asset(
        db_session, monkeypatch, tmp_path,
        tenant_id=tenant["id"], project_id=project["id"], owner_user_id=owner["id"],
    )

    other_tenant, other_user = await _make_tenant_user(client, service_headers, "doc-view-tenant-b")
    r = await client.get(
        f"/api/projects/{project['id']}/assets/{asset.id}/preview",
        headers=_headers(other_tenant["id"], other_user["id"]),
    )
    assert r.status_code in (403, 404)
    r = await client.get(
        f"/api/projects/{project['id']}/assets/{asset.id}/content",
        headers=_headers(other_tenant["id"], other_user["id"]),
    )
    assert r.status_code in (403, 404)


async def test_owner_only_private_asset_is_hidden_from_other_project_members(client, service_headers, db_session, monkeypatch, tmp_path):
    tenant, owner = await _make_tenant_user(client, service_headers, "doc-view-private")
    r = await client.post(
        "/api/projects", json={"name": "P", "description": "", "is_shared": True},
        headers=_headers(tenant["id"], owner["id"], "editor"),
    )
    project = r.json()

    r = await client.post(
        f"/api/tenants/{tenant['id']}/users",
        json={"email": "member@test.com", "display_name": "Member", "role": "editor", "external_id": "ext-member"},
        headers=service_headers,
    )
    member = r.json()

    from app.models.project import ProjectMember

    # Creating the project already added its owner as a member; only the
    # second user's membership needs to be added here.
    db_session.add(ProjectMember(project_id=project["id"], user_id=member["id"], role="member", is_active=True))
    await db_session.commit()

    asset = await _seed_asset(
        db_session, monkeypatch, tmp_path,
        tenant_id=tenant["id"], project_id=project["id"], owner_user_id=owner["id"],
        visibility="private", content=b"top secret",
    )

    r = await client.get(
        f"/api/projects/{project['id']}/assets/{asset.id}/preview",
        headers=_headers(tenant["id"], member["id"]),
    )
    assert r.status_code == 403

    r = await client.get(
        f"/api/projects/{project['id']}/assets/{asset.id}/preview",
        headers=_headers(tenant["id"], owner["id"]),
    )
    assert r.status_code == 200
    assert r.json()["text"] == "top secret"


async def test_storage_path_outside_the_configured_root_is_rejected(client, service_headers, db_session, monkeypatch, tmp_path):
    """Defense in depth: even if a storage_location value somehow points
    outside LOCAL_STORAGE_BASE, the content/preview endpoints must refuse to
    serve it rather than following the path."""
    tenant, owner = await _make_tenant_user(client, service_headers, "doc-view-containment")
    r = await client.post(
        "/api/projects", json={"name": "P", "description": "", "is_shared": True},
        headers=_headers(tenant["id"], owner["id"], "editor"),
    )
    project = r.json()

    import app.routes.project_assets as project_assets_module
    from app.models.project_asset import ProjectAsset

    storage_root = tmp_path / "storage_root"
    storage_root.mkdir()
    monkeypatch.setattr(project_assets_module, "LOCAL_STORAGE_BASE", str(storage_root))

    outside_dir = tmp_path / "outside"
    outside_dir.mkdir()
    escaped_file = outside_dir / "secret.txt"
    escaped_file.write_bytes(b"should never be served")

    asset = ProjectAsset(
        tenant_id=tenant["id"], project_id=project["id"], owner_user_id=owner["id"],
        asset_type="txt", source_type="uploaded_file", title="Doc",
        filename="secret.txt", original_filename="secret.txt",
        content_type="text/plain", file_extension=".txt",
        storage_provider="local", storage_location=str(escaped_file),
        file_size_bytes=len(b"should never be served"), visibility="shared_project", status="uploaded",
    )
    db_session.add(asset)
    await db_session.commit()
    await db_session.refresh(asset)

    r = await client.get(
        f"/api/projects/{project['id']}/assets/{asset.id}/content",
        headers=_headers(tenant["id"], owner["id"]),
    )
    assert r.status_code == 404
